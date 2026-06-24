from io import BytesIO

import pytest

from app.documents.parser import DocumentParseError, parse_uploaded_document
from app.core.settings import get_settings


def test_parse_markdown_uses_utf8_text_path() -> None:
    parsed = parse_uploaded_document(
        filename="guide.md",
        content_type="text/markdown",
        data="# 标题\n\n这是 UTF-8 Markdown 内容。".encode("utf-8"),
        title="知识文档",
        settings=get_settings(),
    )

    assert parsed.title == "知识文档"
    assert parsed.content == "# 标题\n这是 UTF-8 Markdown 内容。"
    assert parsed.doc_type == "markdown"
    assert parsed.metadata["ocr_used"] is False


def test_parse_text_uses_charset_fallback_for_gb18030() -> None:
    settings = get_settings()
    parsed = parse_uploaded_document(
        filename="legacy.txt",
        content_type="text/plain",
        data="中文旧编码内容".encode("gb18030"),
        title=None,
        settings=settings,
    )

    assert parsed.content == "中文旧编码内容"
    assert parsed.doc_type == "text"
    assert parsed.metadata["parser"] == "gb18030"


def test_parse_docx_preserves_paragraph_table_order() -> None:
    from docx import Document

    document = Document()
    document.add_paragraph("开头段落")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "表格左"
    table.rows[0].cells[1].text = "表格右"
    document.add_paragraph("结尾段落")

    buffer = BytesIO()
    document.save(buffer)

    parsed = parse_uploaded_document(
        filename="ordered.docx",
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        data=buffer.getvalue(),
        title=None,
        settings=get_settings(),
    )

    assert parsed.doc_type == "docx"
    assert parsed.content.splitlines() == ["开头段落", "表格左\t表格右", "结尾段落"]


# ---- 结构恢复：docx/pptx 标题 → Markdown `#`（供 chunker 识别章节，新增）----


def _docx_bytes(paragraphs: list[tuple[str, str | None]]) -> bytes:
    """构造 docx：paragraphs 为 (文本, 样式名) 列表，样式名 None 表示普通正文段。"""
    from docx import Document

    document = Document()
    for text, style in paragraphs:
        document.add_paragraph(text, style=style) if style else document.add_paragraph(text)
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def test_parse_docx_heading_styles_become_markdown_prefixes() -> None:
    # Heading 1/2 样式段 → `# `/`## `；正文段不加前缀（结构恢复使 chunker 能切章节）。
    data = _docx_bytes(
        [
            ("数据来源", "Heading 1"),
            ("卫星影像通过开放中心下载。", None),
            ("卫星影像", "Heading 2"),
            ("细节内容若干。", None),
        ]
    )

    parsed = parse_uploaded_document(
        filename="headings.docx",
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        data=data,
        title=None,
        settings=get_settings(),
    )

    lines = parsed.content.splitlines()
    assert "# 数据来源" in lines
    assert "## 卫星影像" in lines
    # 正文段绝不被加 `#`
    assert "卫星影像通过开放中心下载。" in lines
    assert "细节内容若干。" in lines


def test_parse_docx_plain_paragraphs_have_no_prefix() -> None:
    # 无标题样式的纯正文 docx：行为与历史一致，绝不出现 `#`（兼容回归）。
    data = _docx_bytes([("第一段。", None), ("第二段。", None)])

    parsed = parse_uploaded_document(
        filename="plain.docx",
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        data=data,
        title=None,
        settings=get_settings(),
    )

    assert "#" not in parsed.content
    assert parsed.content.splitlines() == ["第一段。", "第二段。"]


def test_parse_pptx_title_placeholder_becomes_heading() -> None:
    # pptx 标题占位符 → `# ` 一级标题；正文文本框不加前缀。
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # 标题+内容版式
    slide.shapes.title.text = "遥感概述"
    body = slide.placeholders[1]
    body.text = "这是正文内容。"

    buffer = BytesIO()
    presentation.save(buffer)

    parsed = parse_uploaded_document(
        filename="deck.pptx",
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        data=buffer.getvalue(),
        title=None,
        settings=get_settings(),
    )

    assert "# 遥感概述" in parsed.content
    assert "这是正文内容。" in parsed.content
    # 正文不应被加标题前缀
    assert "# 这是正文内容。" not in parsed.content


def test_docx_heading_level_handles_missing_or_odd_styles() -> None:
    # 异常兜底：style.name 缺失/非标题样式一律返回 0（当正文），不抛错。
    from app.documents.parser import _docx_heading_level

    class _FakeStyle:
        def __init__(self, name):
            self.name = name

    class _FakePara:
        def __init__(self, name):
            self.style = _FakeStyle(name)

    assert _docx_heading_level(_FakePara("Heading 1")) == 1
    assert _docx_heading_level(_FakePara("Heading 3")) == 3
    assert _docx_heading_level(_FakePara("标题 2")) == 2
    assert _docx_heading_level(_FakePara("Title")) == 1
    assert _docx_heading_level(_FakePara("Normal")) == 0
    assert _docx_heading_level(_FakePara(None)) == 0
    assert _docx_heading_level(_FakePara("")) == 0


def test_parse_pdf_uses_pypdf_when_text_is_sufficient(monkeypatch: pytest.MonkeyPatch) -> None:
    class FakePage:
        def extract_text(self):
            return "PDF 正文内容 " * 20

    class FakeReader:
        def __init__(self, _):
            self.pages = [FakePage()]

    monkeypatch.setattr("pypdf.PdfReader", FakeReader)

    parsed = parse_uploaded_document(
        filename="text.pdf",
        content_type="application/pdf",
        data=b"%PDF fake",
        title=None,
        settings=get_settings().model_copy(update={"document_ocr_min_chars_per_page": 20}),
    )

    assert parsed.doc_type == "pdf"
    assert "PDF 正文内容" in parsed.content
    assert parsed.metadata["parser"] == "pypdf"
    assert parsed.metadata["ocr_used"] is False


def test_parse_pdf_records_failed_pypdf_pages(monkeypatch: pytest.MonkeyPatch) -> None:
    class FailingPage:
        def extract_text(self):
            raise RuntimeError("page stream failed")

    class TextPage:
        def extract_text(self):
            return "PDF 正文内容 " * 20

    class FakeReader:
        def __init__(self, _):
            self.pages = [FailingPage(), TextPage()]

    monkeypatch.setattr("pypdf.PdfReader", FakeReader)

    parsed = parse_uploaded_document(
        filename="partial.pdf",
        content_type="application/pdf",
        data=b"%PDF fake",
        title=None,
        settings=get_settings().model_copy(update={"document_ocr_min_chars_per_page": 20}),
    )

    assert parsed.doc_type == "pdf"
    assert "PDF 正文内容" in parsed.content
    assert parsed.metadata["pypdf_failed_pages"] == [1]
    assert parsed.metadata["pypdf_failed_page_count"] == 1


def test_parse_pdf_reports_ocr_unavailable_when_text_is_empty(monkeypatch: pytest.MonkeyPatch) -> None:
    class FakePage:
        def extract_text(self):
            return ""

    class FakeReader:
        def __init__(self, _):
            self.pages = [FakePage()]

    monkeypatch.setattr("pypdf.PdfReader", FakeReader)
    monkeypatch.setattr("app.documents.parser.shutil.which", lambda _: None)

    with pytest.raises(DocumentParseError) as exc_info:
        parse_uploaded_document(
            filename="scan.pdf",
            content_type="application/pdf",
            data=b"%PDF fake",
            title=None,
            settings=get_settings(),
        )

    assert exc_info.value.code == "OCR_UNAVAILABLE"
    assert "tesseract" in exc_info.value.message


def test_parse_pdf_keeps_partial_text_when_ocr_is_unavailable(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    class FakePage:
        def extract_text(self):
            return "页眉"

    class FakeReader:
        def __init__(self, _):
            self.pages = [FakePage()]

    monkeypatch.setattr("pypdf.PdfReader", FakeReader)
    monkeypatch.setattr("app.documents.parser.shutil.which", lambda _: None)

    parsed = parse_uploaded_document(
        filename="mixed.pdf",
        content_type="application/pdf",
        data=b"%PDF fake",
        title=None,
        settings=get_settings().model_copy(update={"document_ocr_min_chars_per_page": 50}),
    )

    assert parsed.doc_type == "pdf"
    assert parsed.content == "页眉"
    assert parsed.metadata["ocr_fallback_to_pypdf"] is True
    assert parsed.metadata["ocr_error_code"] == "OCR_UNAVAILABLE"


def test_parse_rejects_empty_text_document() -> None:
    with pytest.raises(DocumentParseError) as exc_info:
        parse_uploaded_document(
            filename="empty.txt",
            content_type="text/plain",
            data=b" \n\t ",
            title=None,
            settings=get_settings(),
        )

    assert exc_info.value.code == "DOCUMENT_TEXT_EMPTY"


def test_parse_rejects_file_too_large() -> None:
    with pytest.raises(DocumentParseError) as exc_info:
        parse_uploaded_document(
            filename="large.txt",
            content_type="text/plain",
            data=b"abcd",
            title=None,
            settings=get_settings().model_copy(update={"document_max_file_bytes": 3}),
        )

    assert exc_info.value.code == "FILE_TOO_LARGE"
    assert exc_info.value.status_code == 413


def test_parse_pdf_rejects_too_many_pages(monkeypatch: pytest.MonkeyPatch) -> None:
    class FakePage:
        def extract_text(self):
            return "正文"

    class FakeReader:
        def __init__(self, _):
            self.pages = [FakePage(), FakePage()]

    monkeypatch.setattr("pypdf.PdfReader", FakeReader)

    with pytest.raises(DocumentParseError) as exc_info:
        parse_uploaded_document(
            filename="too-many.pdf",
            content_type="application/pdf",
            data=b"%PDF fake",
            title=None,
            settings=get_settings().model_copy(update={"document_max_pdf_pages": 1}),
        )

    assert exc_info.value.code == "PDF_TOO_MANY_PAGES"
    assert exc_info.value.status_code == 413


def test_parse_pdf_rejects_ocr_when_page_count_exceeds_ocr_limit(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    class FakePage:
        def extract_text(self):
            return ""

    class FakeReader:
        def __init__(self, _):
            self.pages = [FakePage(), FakePage()]

    monkeypatch.setattr("pypdf.PdfReader", FakeReader)

    with pytest.raises(DocumentParseError) as exc_info:
        parse_uploaded_document(
            filename="scan.pdf",
            content_type="application/pdf",
            data=b"%PDF fake",
            title=None,
            settings=get_settings().model_copy(
                update={"document_max_pdf_pages": 10, "document_ocr_max_pages": 1}
            ),
        )

    assert exc_info.value.code == "OCR_TOO_MANY_PAGES"
    assert exc_info.value.status_code == 413
