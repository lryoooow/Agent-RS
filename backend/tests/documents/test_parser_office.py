from __future__ import annotations

from io import BytesIO

import pytest

from app.core.settings import Settings
from app.documents.parser import (
    DocumentParseError,
    SUPPORTED_EXTENSIONS,
    parse_uploaded_document,
)


@pytest.fixture
def settings():
    return Settings(_env_file=None)


def _make_pptx(slides: list[list[str]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for texts in slides:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        tf = box.text_frame
        for i, line in enumerate(texts):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx(sheets: dict[str, list[list]]) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_extensions_registered():
    assert ".pptx" in SUPPORTED_EXTENSIONS
    assert ".xlsx" in SUPPORTED_EXTENSIONS


# ---- PPTX ----

def test_parse_pptx_extracts_text(settings):
    data = _make_pptx([["遥感影像分析", "NDVI 植被指数"], ["第二页内容"]])
    parsed = parse_uploaded_document(
        filename="deck.pptx", content_type=None, data=data, title=None, settings=settings
    )
    assert parsed.doc_type == "pptx"
    assert "遥感影像分析" in parsed.content
    assert "NDVI 植被指数" in parsed.content
    assert "第二页内容" in parsed.content
    assert parsed.metadata["slide_count"] == 2
    assert parsed.metadata["parser"] == "python-pptx"


def test_parse_pptx_empty_raises(settings):
    data = _make_pptx([[]])  # one blank slide, no text
    with pytest.raises(DocumentParseError) as exc:
        parse_uploaded_document(
            filename="empty.pptx", content_type=None, data=data, title=None, settings=settings
        )
    assert exc.value.code == "DOCUMENT_TEXT_EMPTY"


def test_parse_pptx_corrupted_raises(settings):
    with pytest.raises(DocumentParseError) as exc:
        parse_uploaded_document(
            filename="bad.pptx", content_type=None, data=b"not a real pptx", title=None, settings=settings
        )
    assert exc.value.code == "PPTX_PARSE_FAILED"


# ---- XLSX ----

def test_parse_xlsx_multi_sheet(settings):
    data = _make_xlsx({
        "Sheet1": [["name", "value"], ["alpha", 10], ["beta", 20]],
        "数据": [["城市", "温度"], ["北京", 25]],
    })
    parsed = parse_uploaded_document(
        filename="book.xlsx", content_type=None, data=data, title=None, settings=settings
    )
    assert parsed.doc_type == "xlsx"
    assert "# Sheet1" in parsed.content
    assert "# 数据" in parsed.content
    assert "alpha\t10" in parsed.content
    assert "北京\t25" in parsed.content
    assert parsed.metadata["sheet_count"] == 2
    assert parsed.metadata["parser"] == "openpyxl"


def test_parse_xlsx_empty_raises(settings):
    data = _make_xlsx({"Empty": []})
    with pytest.raises(DocumentParseError) as exc:
        parse_uploaded_document(
            filename="empty.xlsx", content_type=None, data=data, title=None, settings=settings
        )
    assert exc.value.code == "DOCUMENT_TEXT_EMPTY"


def test_parse_xlsx_corrupted_raises(settings):
    with pytest.raises(DocumentParseError) as exc:
        parse_uploaded_document(
            filename="bad.xlsx", content_type=None, data=b"garbage bytes", title=None, settings=settings
        )
    assert exc.value.code == "XLSX_PARSE_FAILED"


def test_unsupported_type_message_mentions_new_formats(settings):
    with pytest.raises(DocumentParseError) as exc:
        parse_uploaded_document(
            filename="image.png", content_type=None, data=b"x" * 10, title=None, settings=settings
        )
    assert exc.value.code == "UNSUPPORTED_DOCUMENT_TYPE"
    assert "pptx" in exc.value.message and "xlsx" in exc.value.message
