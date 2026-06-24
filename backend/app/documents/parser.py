from __future__ import annotations

import hashlib
import re
import shutil
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

from app.core.settings import Settings

SUPPORTED_EXTENSIONS = {".txt", ".md", ".markdown", ".pdf", ".docx", ".pptx", ".xlsx"}
TEXT_EXTENSIONS = {".txt", ".md", ".markdown"}
OCR_INSTALL_HINT = (
    "OCR dependencies are unavailable. Install macOS dependencies with "
    "`brew install tesseract tesseract-lang poppler`, or Ubuntu dependencies with "
    "`sudo apt install tesseract-ocr tesseract-ocr-chi-sim tesseract-ocr-chi-tra poppler-utils`."
)


class DocumentParseError(ValueError):
    def __init__(self, code: str, message: str, status_code: int = 400) -> None:
        super().__init__(message)
        self.code = code
        self.message = message
        self.status_code = status_code


@dataclass(frozen=True)
class ParsedDocument:
    title: str
    content: str
    doc_type: str
    metadata: dict[str, Any]


def parse_uploaded_document(
    *,
    filename: str,
    content_type: str | None,
    data: bytes,
    title: str | None,
    settings: Settings,
) -> ParsedDocument:
    if len(data) > settings.document_max_file_bytes:
        raise DocumentParseError("FILE_TOO_LARGE", "File exceeds the configured size limit.", 413)

    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise DocumentParseError(
            "UNSUPPORTED_DOCUMENT_TYPE",
            "Only txt, md, markdown, pdf, docx, pptx, and xlsx files are supported.",
        )

    base_metadata = {
        "filename": filename,
        "mime_type": content_type,
        "file_size": len(data),
        "sha256": hashlib.sha256(data).hexdigest(),
        "extension": extension,
    }

    if extension in TEXT_EXTENSIONS:
        content, parser = _parse_text(data)
        doc_type = "markdown" if extension in {".md", ".markdown"} else "text"
        metadata = {**base_metadata, "parser": parser, "ocr_used": False}
    elif extension == ".pdf":
        content, doc_type, metadata = _parse_pdf(data, base_metadata, settings)
    elif extension == ".pptx":
        content, metadata = _parse_pptx(data, base_metadata)
        doc_type = "pptx"
    elif extension == ".xlsx":
        content, metadata = _parse_xlsx(data, base_metadata)
        doc_type = "xlsx"
    else:
        content, metadata = _parse_docx(data, base_metadata)
        doc_type = "docx"

    content = _normalize_content(content)
    if not content:
        raise DocumentParseError("DOCUMENT_TEXT_EMPTY", "No readable text was extracted from the document.")

    metadata["text_length"] = len(content)
    return ParsedDocument(
        title=(title or Path(filename).stem or "Untitled Document").strip(),
        content=content,
        doc_type=doc_type,
        metadata=metadata,
    )


def _parse_text(data: bytes) -> tuple[str, str]:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk", "big5"):
        try:
            return data.decode(encoding), encoding
        except UnicodeDecodeError:
            continue

    from charset_normalizer import from_bytes

    match = from_bytes(data).best()
    if match is None:
        raise DocumentParseError("DOCUMENT_DECODE_FAILED", "Could not detect text encoding.")
    return str(match), f"charset-normalizer:{match.encoding or 'unknown'}"


def _parse_pdf(
    data: bytes,
    base_metadata: dict[str, Any],
    settings: Settings,
) -> tuple[str, str, dict[str, Any]]:
    from pypdf import PdfReader

    try:
        reader = PdfReader(BytesIO(data))
        page_count = len(reader.pages)
    except Exception as exc:
        raise DocumentParseError("PDF_PARSE_FAILED", "Could not parse the PDF file.") from exc

    if page_count > settings.document_max_pdf_pages:
        raise DocumentParseError("PDF_TOO_MANY_PAGES", "PDF exceeds the configured page limit.", 413)

    page_texts, failed_pages = _extract_pdf_page_texts(reader.pages)
    content = "\n\n".join(text for text in page_texts if text)
    average_chars = len(content.strip()) / max(page_count, 1)
    metadata = {
        **base_metadata,
        "page_count": page_count,
        "parser": "pypdf",
        "ocr_used": False,
        "average_chars_per_page": average_chars,
    }
    if failed_pages:
        metadata["pypdf_failed_pages"] = failed_pages
        metadata["pypdf_failed_page_count"] = len(failed_pages)
    has_pypdf_text = bool(content.strip())
    if has_pypdf_text and average_chars >= settings.document_ocr_min_chars_per_page:
        return content, "pdf", metadata

    try:
        ocr_text = _parse_pdf_with_ocr(data, settings, page_count)
    except DocumentParseError as exc:
        if has_pypdf_text:
            return (
                content,
                "pdf",
                {
                    **metadata,
                    "ocr_attempted": True,
                    "ocr_error_code": exc.code,
                    "ocr_fallback_to_pypdf": True,
                },
            )
        raise

    combined_content = "\n\n".join(text for text in (content, ocr_text) if text.strip())
    if not ocr_text.strip() and has_pypdf_text:
        return (
            content,
            "pdf",
            {
                **metadata,
                "ocr_attempted": True,
                "ocr_text_empty": True,
                "ocr_fallback_to_pypdf": True,
            },
        )

    return (
        combined_content,
        "pdf_ocr",
        {
            **metadata,
            "parser": "pypdf+pdf2image+pytesseract" if has_pypdf_text else "pdf2image+pytesseract",
            "ocr_used": True,
            "ocr_languages": settings.document_ocr_languages,
        },
    )


def _extract_pdf_page_texts(pages: Any) -> tuple[list[str], list[int]]:
    page_texts: list[str] = []
    failed_pages: list[int] = []
    for page_number, page in enumerate(pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            failed_pages.append(page_number)
            continue
        if text:
            page_texts.append(text)
    return page_texts, failed_pages


def _parse_pdf_with_ocr(data: bytes, settings: Settings, page_count: int) -> str:
    if page_count > settings.document_ocr_max_pages:
        raise DocumentParseError(
            "OCR_TOO_MANY_PAGES",
            "PDF requires OCR and exceeds the configured OCR page limit.",
            413,
        )
    if not shutil.which("tesseract") or not shutil.which("pdftoppm"):
        raise DocumentParseError("OCR_UNAVAILABLE", OCR_INSTALL_HINT, 422)

    try:
        from pdf2image import convert_from_bytes
        import pytesseract

        images = convert_from_bytes(data, dpi=200, timeout=settings.document_ocr_timeout_seconds)
        texts = [
            pytesseract.image_to_string(
                image,
                lang=settings.document_ocr_languages,
                timeout=settings.document_ocr_timeout_seconds,
            ).strip()
            for image in images
        ]
    except Exception as exc:
        raise DocumentParseError("OCR_UNAVAILABLE", OCR_INSTALL_HINT, 422) from exc

    return "\n\n".join(text for text in texts if text)


def _docx_heading_level(paragraph: Any) -> int:
    """从 docx 段落样式名解析标题级别（1~6）；非标题/异常返回 0（当正文）。

    兼容英文样式名（Heading 1 / Heading1 / Title）与中文样式名（标题 1）。
    style 缺失或样式名异常时一律降级为 0，绝不因结构恢复破坏既有纯文本路径。
    """
    try:
        name = (paragraph.style.name or "").strip()
    except Exception:
        return 0
    if not name:
        return 0
    lowered = name.lower()
    # 文档主标题（Title / 标题）作一级。
    if lowered == "title" or name == "标题":
        return 1
    match = re.search(r"(?:heading|标题)\s*([1-6])", lowered if "heading" in lowered else name)
    if match:
        return int(match.group(1))
    return 0


def _parse_docx(data: bytes, base_metadata: dict[str, Any]) -> tuple[str, dict[str, Any]]:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    try:
        document = Document(BytesIO(data))
    except Exception as exc:
        raise DocumentParseError("DOCX_PARSE_FAILED", "Could not parse the docx file.") from exc

    parts: list[str] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if text:
                # 结构恢复：Word 的 Heading 样式段落 → Markdown `#` 前缀，供下游 chunker
                # 识别章节并生成面包屑。原先只取 .text、丢弃样式，docx 章节信息全失（覆盖率<10%）。
                level = _docx_heading_level(paragraph)
                parts.append(f"{'#' * level} {text}" if level else text)
        elif isinstance(child, CT_Tbl):
            table = Table(child, document)
            rows = [
                "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                for row in table.rows
            ]
            table_text = "\n".join(row for row in rows if row)
            if table_text:
                parts.append(table_text)

    return "\n\n".join(parts), {**base_metadata, "parser": "python-docx", "ocr_used": False}


def _pptx_shape_is_title(shape: Any) -> bool:
    """判断 pptx 形状是否为标题占位符（TITLE / CENTER_TITLE）。

    非占位符形状访问 placeholder_format.type 会抛错或返回 None，故 try 兜底为 False，
    不因结构恢复破坏既有文本路径。
    """
    try:
        if not shape.is_placeholder:
            return False
        ph_type = shape.placeholder_format.type
    except Exception:
        return False
    if ph_type is None:
        return False
    # PP_PLACEHOLDER.TITLE=13, CENTER_TITLE=0；用名字判断避免依赖枚举导入。
    return str(ph_type).upper().find("TITLE") != -1


def _parse_pptx(data: bytes, base_metadata: dict[str, Any]) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    try:
        presentation = Presentation(BytesIO(data))
    except Exception as exc:
        raise DocumentParseError("PPTX_PARSE_FAILED", "Could not parse the pptx file.") from exc

    slide_texts: list[str] = []
    slide_count = 0
    for slide in presentation.slides:
        slide_count += 1
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    para.text.strip()
                    for para in shape.text_frame.paragraphs
                    if para.text.strip()
                )
                if text:
                    # 结构恢复：标题占位符 → Markdown 一级标题，供下游 chunker 识别章节。
                    parts.append(f"# {text}" if _pptx_shape_is_title(shape) else text)
            if shape.has_table:
                rows = [
                    "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    for row in shape.table.rows
                ]
                table_text = "\n".join(row for row in rows if row)
                if table_text:
                    parts.append(table_text)
        if parts:
            slide_texts.append("\n".join(parts))

    content = "\n\n".join(slide_texts)
    return content, {
        **base_metadata,
        "parser": "python-pptx",
        "ocr_used": False,
        "slide_count": slide_count,
    }


def _parse_xlsx(data: bytes, base_metadata: dict[str, Any]) -> tuple[str, dict[str, Any]]:
    from openpyxl import load_workbook

    try:
        # read_only 流式读取省内存；data_only 取缓存值而非公式串。
        workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        raise DocumentParseError("XLSX_PARSE_FAILED", "Could not parse the xlsx file.") from exc

    try:
        sheet_blocks: list[str] = []
        sheet_count = 0
        for worksheet in workbook.worksheets:
            sheet_count += 1
            row_texts: list[str] = []
            for row in worksheet.iter_rows(values_only=True):
                cells = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if cells:
                    row_texts.append("\t".join(cells))
            if row_texts:
                sheet_blocks.append(f"# {worksheet.title}\n" + "\n".join(row_texts))
    finally:
        workbook.close()

    content = "\n\n".join(sheet_blocks)
    return content, {
        **base_metadata,
        "parser": "openpyxl",
        "ocr_used": False,
        "sheet_count": sheet_count,
    }


def _normalize_content(content: str) -> str:
    return "\n".join(line.strip() for line in content.splitlines() if line.strip()).strip()
